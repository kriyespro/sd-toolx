"""PDF to PowerPoint — one slide per page (image of each page)."""
import io

import fitz
from pptx import Presentation
from pptx.util import Emu, Pt

from features.tools.models import ToolJob
from features.tools.services import file_handler
from features.tools.services.exceptions import ToolError


def run(job: ToolJob) -> dict:
    src_rel = job.input_files[0]
    src = file_handler.to_absolute(src_rel)
    out = file_handler.output_path(job, "converted.pptx")

    doc = fitz.open(str(src))
    if doc.page_count == 0:
        doc.close()
        raise ToolError("PDF has no pages.")

    prs = Presentation()
    first_page = doc[0]
    width_pt = first_page.rect.width
    height_pt = first_page.rect.height
    prs.slide_width = Emu(int(width_pt * 12700))
    prs.slide_height = Emu(int(height_pt * 12700))

    blank_layout = prs.slide_layouts[6]

    for page in doc:
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_bytes = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_bytes),
            0, 0,
            prs.slide_width,
            prs.slide_height,
        )

    slide_count = doc.page_count
    doc.close()
    prs.save(str(out))
    out_rel = file_handler.to_relative(out)
    return {
        "output_file": out_rel,
        "compression_stat": {"slides": slide_count},
        "preview_file": "",
    }
