"""PowerPoint to PDF via LibreOffice (fallback: python-pptx + PyMuPDF)."""
import io
import shutil
import subprocess
import tempfile
from pathlib import Path

import fitz
from PIL import Image

from features.tools.models import ToolJob
from features.tools.services import file_handler
from features.tools.services.exceptions import ToolError
from features.tools.services.pdf_utils import render_preview


def _find_soffice():
    for cmd in ("soffice", "libreoffice", "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        if cmd.startswith("/"):
            if Path(cmd).exists():
                return cmd
        else:
            found = shutil.which(cmd)
            if found:
                return found
    return None


def _ppt_to_pdf_fallback(src: Path, out: Path, job: ToolJob) -> dict:
    """Convert PPTX to PDF by rendering each slide's text with PyMuPDF."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as exc:
        raise ToolError("python-pptx not installed.") from exc

    try:
        prs = Presentation(str(src))
    except Exception as exc:
        raise ToolError(f"Cannot read PowerPoint file: {exc}") from exc

    doc = fitz.open()
    for i, slide in enumerate(prs.slides):
        page = doc.new_page(width=842, height=595)
        y = 60

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                if y > 560:
                    break
                is_title = shape.shape_type == 13 or (hasattr(shape, 'placeholder_format') and
                           shape.placeholder_format and shape.placeholder_format.idx == 0)
                fontsize = 18 if is_title else 11
                page.insert_text(fitz.Point(40, y), text[:120], fontsize=fontsize)
                y += fontsize + 6

        page.insert_text(fitz.Point(800, 580), str(i + 1), fontsize=8, color=(0.5, 0.5, 0.5))

    doc.save(str(out), garbage=3, deflate=True)
    doc.close()
    out_rel = file_handler.to_relative(out)
    return {
        "output_file": out_rel,
        "compression_stat": {"converter": "pptx-fallback", "slides": len(prs.slides)},
        "preview_file": render_preview(job, out_rel),
    }


def run(job: ToolJob) -> dict:
    src_rel = job.input_files[0]
    src = file_handler.to_absolute(src_rel)
    out = file_handler.output_path(job, "converted.pdf")
    soffice = _find_soffice()

    if not soffice:
        return _ppt_to_pdf_fallback(src, out, job)

    with tempfile.TemporaryDirectory() as tmp:
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(src)],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            raise ToolError("Conversion failed. Ensure the PowerPoint file is valid.")
        pdfs = list(Path(tmp).glob("*.pdf"))
        if not pdfs:
            raise ToolError("No PDF output produced.")
        shutil.copy(pdfs[0], out)

    out_rel = file_handler.to_relative(out)
    return {
        "output_file": out_rel,
        "compression_stat": {"converter": "libreoffice"},
        "preview_file": render_preview(job, out_rel),
    }
